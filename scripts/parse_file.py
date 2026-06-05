#!/usr/bin/env python3
"""
Convert uploaded files to Markdown.

Primary strategy:
1. Use markitdown to normalize every supported file into Markdown.
2. Fallback to format-specific extractors if markitdown is unavailable.
3. When a subject is provided, save each file under storage/<subject>/markdown格式/.
"""
from __future__ import annotations

import argparse
import io
import sys
from pathlib import Path

try:
    from .utils import save_markdown_copy, save_parsed_content
except ImportError:
    sys.path.append(str(Path(__file__).resolve().parent))
    from utils import save_markdown_copy, save_parsed_content


def _fix_encoding() -> None:
    if sys.platform == "win32":
        try:
            sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding="utf-8")
            sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding="utf-8")
        except Exception:
            pass


def parse_with_markitdown(file_path: Path) -> tuple[str | None, str | None]:
    try:
        from markitdown import MarkItDown
    except ImportError:
        return None, "missing markitdown"

    try:
        converter = MarkItDown(enable_plugins=False)
        result = converter.convert(str(file_path))
        content = (
            getattr(result, "text_content", None)
            or getattr(result, "markdown", None)
            or getattr(result, "content", None)
        )
        if content is None:
            content = str(result)
        return content.strip(), None
    except Exception as exc:
        return None, f"markitdown failed: {exc}"


def parse_pptx(file_path: Path) -> tuple[str | None, str | None]:
    try:
        from pptx import Presentation
    except ImportError:
        return None, "缺少 python-pptx。请运行: pip install python-pptx"

    presentation = Presentation(file_path)
    slides = []
    for index, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"## 第{index}页\n\n" + "\n\n".join(texts))
    return "\n\n".join(slides), None


def parse_pdf(file_path: Path) -> tuple[str | None, str | None]:
    try:
        import pdfplumber

        sections = []
        with pdfplumber.open(file_path) as pdf:
            for index, page in enumerate(pdf.pages, start=1):
                text = page.extract_text()
                if text and text.strip():
                    sections.append(f"## 第{index}页\n\n{text.strip()}")
        return "\n\n".join(sections), None
    except ImportError:
        pass
    except Exception as exc:
        return None, f"pdfplumber failed: {exc}"

    try:
        from PyPDF2 import PdfReader

        sections = []
        reader = PdfReader(str(file_path))
        for index, page in enumerate(reader.pages, start=1):
            text = page.extract_text()
            if text and text.strip():
                sections.append(f"## 第{index}页\n\n{text.strip()}")
        return "\n\n".join(sections), None
    except ImportError:
        return None, "缺少 pdfplumber 或 PyPDF2。请运行: pip install pdfplumber PyPDF2"
    except Exception as exc:
        return None, f"PyPDF2 failed: {exc}"


def parse_word(file_path: Path) -> tuple[str | None, str | None]:
    try:
        from docx import Document
    except ImportError:
        return None, "缺少 python-docx。请运行: pip install python-docx"

    document = Document(file_path)
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n\n".join(paragraphs), None


def parse_image(file_path: Path) -> tuple[str | None, str | None]:
    try:
        from PIL import Image
        import pytesseract
    except ImportError:
        return None, "缺少 pytesseract 或 Pillow。请运行: pip install pytesseract Pillow"

    image = Image.open(file_path)
    for lang in ("chi_sim+eng", "chi_sim", "eng"):
        try:
            text = pytesseract.image_to_string(image, lang=lang)
            if text.strip():
                return text.strip(), None
        except Exception:
            continue
    return None, "OCR 未识别到可用文本"


def parse_text(file_path: Path) -> tuple[str | None, str | None]:
    for encoding in ("utf-8", "gbk", "utf-16"):
        try:
            return file_path.read_text(encoding=encoding), None
        except UnicodeDecodeError:
            continue
    return None, f"无法解码文本文件: {file_path}"


def parse_file(file_path: str | Path) -> tuple[str | None, str | None]:
    path = Path(file_path)
    if not path.exists():
        return None, f"文件不存在: {path}"

    markdown, error = parse_with_markitdown(path)
    if markdown:
        return markdown, None

    suffix = path.suffix.lower()
    parsers = {
        ".pptx": parse_pptx,
        ".pdf": parse_pdf,
        ".doc": parse_word,
        ".docx": parse_word,
        ".png": parse_image,
        ".jpg": parse_image,
        ".jpeg": parse_image,
        ".bmp": parse_image,
        ".gif": parse_image,
        ".webp": parse_image,
        ".txt": parse_text,
        ".md": parse_text,
    }

    parser = parsers.get(suffix)
    if not parser:
        return None, f"不支持的文件格式: {suffix}"
    return parser(path)


def parse_multiple(files: list[str], subject_name: str | None = None) -> tuple[str, list[str], list[str]]:
    merged_sections: list[str] = []
    saved_paths: list[str] = []
    errors: list[str] = []

    for file_name in files:
        path = Path(file_name)
        markdown, error = parse_file(path)
        if error or not markdown:
            errors.append(f"{path}: {error or '未提取到内容'}")
            continue

        section = f"# {path.stem}\n\n{markdown.strip()}"
        merged_sections.append(section)
        if subject_name:
            saved_paths.append(save_markdown_copy(subject_name, path.name, section))

    merged_markdown = "\n\n---\n\n".join(merged_sections)
    if subject_name and merged_markdown:
        save_parsed_content(subject_name, merged_markdown)

    return merged_markdown, saved_paths, errors


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="Convert review files to Markdown.")
    parser.add_argument("files", nargs="+", help="One or more files to parse")
    parser.add_argument("--subject", help="Subject name. If provided, save outputs into storage/<subject>/markdown格式/")
    parser.add_argument("--stdout", action="store_true", help="Always print merged markdown to stdout")
    return parser


def main() -> None:
    _fix_encoding()
    parser = build_parser()
    args = parser.parse_args()

    markdown, saved_paths, errors = parse_multiple(args.files, args.subject)

    for error in errors:
        print(f"[错误] {error}", file=sys.stderr)

    if args.subject:
        print(f"科目: {args.subject}")
        print(f"已保存 Markdown 文件: {len(saved_paths)}")
        for saved_path in saved_paths:
            print(f"  - {saved_path}")

    if markdown and (args.stdout or not args.subject):
        print(markdown)

    if not markdown:
        print("未能提取到任何内容", file=sys.stderr)
        sys.exit(1)


if __name__ == "__main__":
    main()
